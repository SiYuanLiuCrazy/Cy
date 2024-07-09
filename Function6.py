from pptx import Presentation
import os

def merge_pptx(input_dir, output_dir):
    # 初始化合并后的演示文稿
    merged_presentation = Presentation()

    # 获取输入目录中的文件列表
    files_to_merge = os.listdir(input_dir)
    files_to_merge.sort()  # 确保文件按照指定顺序合并

    # 将每个 PPTX 文件合并到 merged_presentation 中
    for filename in files_to_merge:
        if filename.endswith(".pptx") or filename.endswith(".ppt"):
            file_path = os.path.join(input_dir, filename)
            presentation_to_merge = Presentation(file_path)
            for slide in presentation_to_merge.slides:
                # 复制幻灯片到合并的演示文稿中
                merged_slide = merged_presentation.slides.add_slide(slide.slide_layout)
                for shape in slide.shapes:
                    merged_slide.shapes._spTree.append(shape.element)

    # 用户输入合并后的文件名
    output_filename = input("请输入合并后的文件名（包括扩展名）：")

    # 保存合并后的演示文稿
    output_path = os.path.join(output_dir, output_filename)
    merged_presentation.save(output_path)
    print(f"合并后的演示文稿已保存至：{output_path}")

if __name__ == "__main__":
    input_directory = r'E:\PPT_ReadyToMerge'
    output_directory = r'E:\PPT_Merge'
    merge_pptx(input_directory, output_directory)
