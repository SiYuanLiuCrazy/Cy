from pptx import Presentation
import os
from io import BytesIO

def copy_slide(prs, slide):
    # 创建一个新的幻灯片，使用默认布局
    new_slide_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(new_slide_layout)

    # 复制当前幻灯片的所有形状到新的幻灯片中
    for shape in slide.shapes:
        if shape.shape_type == 1:  # 处理文本框
            new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text = shape.text
        elif shape.shape_type == 13:  # 处理图片
            image_stream = BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
        else:
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text = shape.text_frame.text
            else:
                new_slide.shapes._spTree.insert_element_before(shape._element, 'p:extLst')

    # 复制背景
    new_slide_background = new_slide.background
    slide_background = slide.background
    if slide_background.fill.type == 'solid':
        new_slide_background.fill.solid()
        new_slide_background.fill.fore_color.rgb = slide_background.fill.fore_color.rgb
    elif slide_background.fill.type == 'patterned':
        new_slide_background.fill.patterned()
        new_slide_background.fill.fore_color.rgb = slide_background.fill.fore_color.rgb

    return new_slide

def split_pptx(input_path, output_dir):
    # 加载PPTX文件
    presentation = Presentation(input_path)
    base_name = os.path.splitext(os.path.basename(input_path))[0]

    # 创建保存单独PPTX文件的目录
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # 遍历每一页
    for i, slide in enumerate(presentation.slides):
        # 创建一个新的演示文稿
        new_presentation = Presentation()
        # 复制当前幻灯片到新的演示文稿中
        copy_slide(new_presentation, slide)

        # 保存单独的PPTX文件
        new_file_path = os.path.join(output_dir, f"{base_name}_slide_{i+1}.pptx")
        new_presentation.save(new_file_path)
        print(f"Saved {new_file_path}")

# 遍历E:\PPT_New文件夹中的所有PPTX文件
input_directory = r'E:\PPT_New'
output_directory = r'E:\PPT_Signal_Page'

for filename in os.listdir(input_directory):
    if filename.endswith('.pptx'):
        input_pptx_path = os.path.join(input_directory, filename)
        split_pptx(input_pptx_path, output_directory)
