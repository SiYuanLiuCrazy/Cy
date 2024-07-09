import os
from pptx import Presentation

def split_pptx(file_path, output_dir):
    # 读取PPTX文件
    presentation = Presentation(file_path)
    num_slides = len(presentation.slides)
    print(f"Processing {file_path} with {num_slides} slides")

    # 获取原始文件名，不包含扩展名
    base_filename = os.path.splitext(os.path.basename(file_path))[0]

    # 创建以原文件名命名的子目录
    specific_output_dir = os.path.join(output_dir, base_filename)
    if not os.path.exists(specific_output_dir):
        os.makedirs(specific_output_dir)

    for i in range(num_slides):
        # 创建一个新的PPTX文件副本
        single_slide_presentation = Presentation(file_path)
        
        # 获取所有幻灯片的索引
        slide_indexes = list(range(len(single_slide_presentation.slides)))
        
        # 删除除了第i个幻灯片以外的所有幻灯片
        for index in sorted(slide_indexes, reverse=True):
            if index != i:
                # 直接从XML中删除幻灯片
                slide_id = single_slide_presentation.slides._sldIdLst[index].rId
                single_slide_presentation.part.drop_rel(slide_id)
                del single_slide_presentation.slides._sldIdLst[index]
        
        # 生成新的文件名，包含原始文件名和页数
        single_slide_filename = os.path.join(specific_output_dir, f"{base_filename}_slide_{i+1}.pptx")
        single_slide_presentation.save(single_slide_filename)
        print(f"Saved slide {i+1} as {single_slide_filename}")

def process_pptx_folder(input_dir, output_dir):
    # 遍历输入目录中的所有PPTX文件
    for filename in os.listdir(input_dir):
        if filename.endswith(".pptx"):
            file_path = os.path.join(input_dir, filename)
            split_pptx(file_path, output_dir)

input_dir = "E:\\PPT_Home"
output_dir = "E:\\PPT_Signal_Page"

process_pptx_folder(input_dir, output_dir)